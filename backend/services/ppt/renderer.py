"""Render a validated PptAgent outline into a polished 16:9 PowerPoint deck."""

from __future__ import annotations

import json
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5
_LAYOUT_BLANK = 6

_BG = RGBColor(0xF3, 0xF7, 0xF6)
_PAPER = RGBColor(0xFF, 0xFF, 0xFF)
_INK = RGBColor(0x14, 0x2F, 0x2E)
_MUTED = RGBColor(0x5D, 0x73, 0x71)
_TEAL = RGBColor(0x0F, 0x85, 0x88)
_TEAL_DARK = RGBColor(0x0A, 0x4F, 0x52)
_MINT = RGBColor(0xD9, 0xEB, 0xE9)
_MINT_SOFT = RGBColor(0xEA, 0xF3, 0xF1)
_AMBER = RGBColor(0xD2, 0x91, 0x3D)
_CODE_BG = RGBColor(0x0E, 0x1D, 0x1C)
_CODE_FG = RGBColor(0xE8, 0xF2, 0xEF)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_BODY_FONT = "Microsoft YaHei"
_CODE_FONT = "Consolas"


def _set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _shape(
    slide,
    shape_type: MSO_SHAPE,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius: bool = False,
):
    actual_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type
    item = slide.shapes.add_shape(
        actual_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    item.fill.solid()
    item.fill.fore_color.rgb = fill
    item.line.color.rgb = line or fill
    return item


def _text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: RGBColor = _INK,
    bold: bool = False,
    font: str = _BODY_FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _set_notes(slide, notes: str) -> None:
    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()


def _eyebrow(slide, label: str, *, dark: bool = False) -> None:
    color = _MINT if dark else _TEAL
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 0.48, 0.08, 0.28, fill=_AMBER)
    _text(
        slide,
        label.upper(),
        0.94,
        0.42,
        4.8,
        0.38,
        size=11,
        color=color,
        bold=True,
    )


def _title(slide, title: str, *, dark: bool = False) -> None:
    _eyebrow(slide, "ALGOPILOT · COURSE DECK", dark=dark)
    _text(
        slide,
        title,
        0.72,
        0.92,
        11.8,
        0.72,
        size=35,
        color=_WHITE if dark else _INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _footer(slide, deck_title: str, index: int, total: int, *, dark: bool = False) -> None:
    line_color = RGBColor(0x2E, 0x66, 0x64) if dark else RGBColor(0xCE, 0xDD, 0xDA)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 7.08, 11.9, 0.012, fill=line_color)
    _text(
        slide,
        deck_title[:34],
        0.72,
        7.12,
        6.5,
        0.2,
        size=9,
        color=_MINT if dark else _MUTED,
    )
    _text(
        slide,
        f"{index:02d} / {total:02d}",
        11.3,
        7.1,
        1.3,
        0.22,
        size=10,
        color=_MINT if dark else _TEAL,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _render_cover(slide, data: dict[str, Any], deck_title: str) -> None:
    _set_background(slide, _TEAL_DARK)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, fill=_AMBER)
    _shape(slide, MSO_SHAPE.OVAL, 10.6, -1.0, 4.2, 4.2, fill=_TEAL, line=_TEAL)
    _shape(slide, MSO_SHAPE.OVAL, 11.35, 4.9, 2.6, 2.6, fill=_TEAL, line=_TEAL)
    _text(
        slide,
        "ALGOPILOT · 个性化课程讲义",
        0.88,
        0.7,
        5.8,
        0.4,
        size=14,
        color=_MINT,
        bold=True,
    )
    title = str(data.get("title") or deck_title or "课程讲义").strip()
    _text(
        slide,
        title,
        0.88,
        2.0,
        10.8,
        1.65,
        size=52,
        color=_WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    subtitle = str(data.get("subtitle") or "从概念理解到可验证实现").strip()
    _text(slide, subtitle, 0.92, 3.9, 9.6, 0.62, size=22, color=_MINT)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.88, 5.55, 5.4, 0.68, fill=_TEAL, line=_TEAL, radius=True)
    _text(
        slide,
        "理解方法  ·  跟踪过程  ·  验证边界",
        1.12,
        5.71,
        4.9,
        0.34,
        size=16,
        color=_WHITE,
        bold=True,
    )
    _set_notes(slide, str(data.get("notes") or ""))


def _render_agenda(slide, data: dict[str, Any]) -> None:
    _set_background(slide, _BG)
    _title(slide, str(data.get("title") or "学习路径"))
    bullets = [str(item).strip() for item in data.get("bullets") or [] if str(item).strip()][:5]
    top = 2.05
    for idx, item in enumerate(bullets, start=1):
        _text(slide, f"{idx:02d}", 0.84, top, 0.72, 0.48, size=22, color=_TEAL, bold=True)
        _text(slide, item, 1.72, top - 0.02, 9.8, 0.5, size=22, color=_INK, bold=True)
        _shape(slide, MSO_SHAPE.RECTANGLE, 1.72, top + 0.58, 10.45, 0.012, fill=RGBColor(0xD3, 0xE0, 0xDD))
        top += 0.9
    _set_notes(slide, str(data.get("notes") or ""))


def _render_content(slide, data: dict[str, Any]) -> None:
    _set_background(slide, _BG)
    _title(slide, str(data.get("title") or "核心要点"))
    bullets = [str(item).strip() for item in data.get("bullets") or [] if str(item).strip()][:5]
    if not bullets:
        bullets = ["本页内容待补充"]
    row_height = min(0.94, 4.65 / len(bullets))
    top = 1.95
    for idx, item in enumerate(bullets, start=1):
        _shape(slide, MSO_SHAPE.OVAL, 0.82, top + 0.05, 0.44, 0.44, fill=_TEAL, line=_TEAL)
        _text(
            slide,
            str(idx),
            0.82,
            top + 0.05,
            0.44,
            0.44,
            size=13,
            color=_WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        _text(
            slide,
            item,
            1.52,
            top,
            10.65,
            row_height - 0.05,
            size=21,
            color=_INK,
            bold=idx == 1,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if idx < len(bullets):
            _shape(slide, MSO_SHAPE.RECTANGLE, 1.52, top + row_height - 0.08, 10.45, 0.012, fill=RGBColor(0xD5, 0xE2, 0xDF))
        top += row_height
    _shape(slide, MSO_SHAPE.RECTANGLE, 12.82, 1.92, 0.12, min(4.75, len(bullets) * row_height), fill=_MINT)
    _set_notes(slide, str(data.get("notes") or ""))


def _render_code(slide, data: dict[str, Any]) -> None:
    _set_background(slide, _TEAL_DARK)
    _title(slide, str(data.get("title") or "关键实现"), dark=True)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 1.82, 11.9, 4.93, fill=_CODE_BG, line=RGBColor(0x2B, 0x50, 0x4D), radius=True)
    _shape(slide, MSO_SHAPE.OVAL, 1.02, 2.08, 0.12, 0.12, fill=RGBColor(0xE6, 0x6A, 0x62))
    _shape(slide, MSO_SHAPE.OVAL, 1.25, 2.08, 0.12, 0.12, fill=RGBColor(0xE0, 0xB1, 0x48))
    _shape(slide, MSO_SHAPE.OVAL, 1.48, 2.08, 0.12, 0.12, fill=RGBColor(0x67, 0xB9, 0x7A))
    code = str(data.get("code") or "# 暂无代码示例").rstrip()
    line_count = max(1, len(code.splitlines()))
    size = 18 if line_count <= 10 else 16
    _text(
        slide,
        code,
        1.02,
        2.45,
        11.0,
        3.95,
        size=size,
        color=_CODE_FG,
        font=_CODE_FONT,
    )
    _set_notes(slide, str(data.get("notes") or ""))


def _render_closing(slide, data: dict[str, Any]) -> None:
    _set_background(slide, _TEAL_DARK)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, fill=_AMBER)
    _text(slide, "带走三件事", 0.88, 0.78, 3.0, 0.38, size=14, color=_MINT, bold=True)
    _text(
        slide,
        str(data.get("title") or "总结与行动"),
        0.88,
        1.35,
        10.8,
        0.9,
        size=42,
        color=_WHITE,
        bold=True,
    )
    bullets = [str(item).strip() for item in data.get("bullets") or [] if str(item).strip()][:4]
    top = 2.85
    for idx, item in enumerate(bullets, start=1):
        _text(slide, f"0{idx}", 0.96, top, 0.62, 0.45, size=17, color=_AMBER, bold=True)
        _text(slide, item, 1.82, top - 0.02, 9.8, 0.52, size=22, color=_WHITE, bold=idx == 1)
        top += 0.84
    _text(slide, "先定义，再转移；先验证，再优化。", 0.92, 6.35, 8.2, 0.42, size=16, color=_MINT)
    _set_notes(slide, str(data.get("notes") or ""))


_RENDERERS = {
    "agenda": _render_agenda,
    "content": _render_content,
    "code": _render_code,
    "closing": _render_closing,
}


def render_pptx_bytes(outline: dict[str, Any]) -> bytes:
    """Render a normalized outline to a PowerPoint byte stream."""
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)
    deck_title = str(outline.get("title") or "数据结构与算法课程讲义").strip()
    slides = outline.get("slides")
    slide_items = [item for item in slides if isinstance(item, dict)] if isinstance(slides, list) else []
    total = len(slide_items)
    for index, data in enumerate(slide_items, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK])
        layout = str(data.get("layout") or "content").strip().lower()
        if layout == "cover":
            _render_cover(slide, data, deck_title)
        else:
            renderer = _RENDERERS.get(layout, _render_content)
            renderer(slide, data)
            _footer(slide, deck_title, index, total, dark=layout in {"code", "closing"})
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def render_pptx_bytes_from_json(raw_json: str) -> bytes:
    """Parse stored PptAgent JSON and render it to a PowerPoint byte stream."""
    try:
        outline = json.loads(raw_json)
    except (json.JSONDecodeError, TypeError) as exc:
        raise ValueError(f"PPT 大纲 JSON 解析失败：{exc}") from exc
    if not isinstance(outline, dict):
        raise ValueError("PPT 大纲必须是 JSON 对象")
    return render_pptx_bytes(outline)
